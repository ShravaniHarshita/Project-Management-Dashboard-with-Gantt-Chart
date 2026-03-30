from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r"c:\Users\NSVLJ\Downloads\sampleppt (2).pptx")
print(f"Slide width: {prs.slide_width} ({prs.slide_width / 914400:.1f}in)")
print(f"Slide height: {prs.slide_height} ({prs.slide_height / 914400:.1f}in)")
print(f"Slide count: {len(prs.slides)}")
print(f"Layouts: {len(prs.slide_layouts)}")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")
print()

for sn, slide in enumerate(prs.slides, 1):
    print(f"--- SLIDE {sn} (layout: {slide.slide_layout.name}) ---")
    for shape in slide.shapes:
        txt = shape.text[:150].replace("\n", " | ") if hasattr(shape, "text") and shape.text else ""
        print(f"  {shape.shape_type}: L={shape.left} T={shape.top} W={shape.width} H={shape.height}")
        if txt:
            print(f"    TEXT: {txt}")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs[:4]:
                for r in p.runs[:2]:
                    c = None
                    try:
                        c = r.font.color.rgb
                    except:
                        pass
                    print(f"    run: bold={r.font.bold} size={r.font.size} color={c} name={r.font.name} text=\"{r.text[:80]}\"")
        if shape.has_table:
            t = shape.table
            print(f"    TABLE: {len(t.rows)}x{len(t.columns)}")
            for ri, row in enumerate(t.rows):
                cells = [row.cells[ci].text[:40] for ci in range(len(t.columns))]
                print(f"      row{ri}: {cells}")
                if ri > 5:
                    print("      ...")
                    break
    print()
