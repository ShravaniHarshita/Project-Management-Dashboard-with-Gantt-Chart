from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import json

prs = Presentation(r"c:\Users\NSVLJ\Downloads\sampleppt (1).pptx")

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slide count: {len(prs.slides)}")
print(f"Slide layouts available: {len(prs.slide_layouts)}")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: '{layout.name}' - placeholders: {len(layout.placeholders)}")

print("\n" + "="*80)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_num} (layout: '{slide.slide_layout.name}')")
    print(f"{'='*80}")
    
    for shape in slide.shapes:
        print(f"\n  Shape: {shape.shape_type}, Name: '{shape.name}'")
        print(f"    Position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
        
        if hasattr(shape, "text") and shape.text:
            print(f"    Text: '{shape.text[:200]}'")
        
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_info = {
                        'text': run.text[:100],
                        'bold': run.font.bold,
                        'size': str(run.font.size) if run.font.size else None,
                        'color': str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None,
                        'font_name': run.font.name,
                    }
                    print(f"    Run: {font_info}")
        
        if shape.has_table:
            table = shape.table
            print(f"    Table: {len(table.rows)} rows x {len(table.columns)} cols")
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text.strip():
                        print(f"      [{row_idx},{col_idx}]: '{cell.text[:80]}'")
