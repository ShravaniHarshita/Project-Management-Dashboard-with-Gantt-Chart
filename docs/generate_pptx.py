"""
Generate Project Management Dashboard PPTX
Styled to match the sample college presentation format
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Use sample as template ───
SAMPLE_PATH = r"c:\Users\NSVLJ\Downloads\sampleppt (1).pptx"
OUTPUT_PATH = r"c:\Projects\project-management-dashboard\docs\Project_Management_Dashboard_Presentation.pptx"

# Load sample to extract slide dimensions, then create fresh presentation
sample_prs = Presentation(SAMPLE_PATH)
prs = Presentation()
prs.slide_width = sample_prs.slide_width
prs.slide_height = sample_prs.slide_height

# Get layout references
blank_layout = prs.slide_layouts[6]  # Blank layout

# ─── Color Palette ───
DARK_BLUE    = RGBColor(0x1E, 0x3A, 0x5F)
MEDIUM_BLUE  = RGBColor(0x3C, 0x6C, 0xA8)
LIGHT_BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY  = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
BG_BLUE      = RGBColor(0xE8, 0xF0, 0xFE)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_shape_with_fill(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=DARK_GRAY, font_name='Calibri', bullet_char='•'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{bullet_char} {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = font_name
    return txBox


def add_header_bar(slide, title_text):
    """Add a blue header bar at top of slide with white title"""
    bar = add_shape_with_fill(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.1), DARK_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(15), Inches(0.8),
                title_text, font_size=32, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name='Calibri')
    # Accent line under header
    add_shape_with_fill(slide, Emu(0), Inches(1.1), SLIDE_W, Inches(0.06), MEDIUM_BLUE)


def add_section_title(slide, text, left, top, width=Inches(10)):
    add_textbox(slide, left, top, width, Inches(0.5), text,
                font_size=24, font_color=DARK_BLUE, bold=True, font_name='Calibri')
    # underline
    add_shape_with_fill(slide, left, top + Inches(0.45), Inches(3), Inches(0.04), MEDIUM_BLUE)


def add_table(slide, left, top, width, height, rows_data, col_widths=None):
    """rows_data = list of lists of strings. First row is header."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 1
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape


def add_info_card(slide, left, top, width, height, title, body_items, title_color=DARK_BLUE):
    """Add a rounded card with title and bullet body"""
    card = add_shape_with_fill(slide, left, top, width, height, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.shadow.inherit = False
    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
                title, font_size=16, font_color=title_color, bold=True)
    # Bullets
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4),
                    height - Inches(0.6), body_items, font_size=12, font_color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full blue background
add_shape_with_fill(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, DARK_BLUE)
# Accent bar
add_shape_with_fill(slide, Emu(0), Inches(3.0), SLIDE_W, Inches(0.08), LIGHT_BLUE)
add_shape_with_fill(slide, Emu(0), Inches(4.7), SLIDE_W, Inches(0.08), LIGHT_BLUE)

add_textbox(slide, Inches(1), Inches(0.6), Inches(16), Inches(0.6),
            'RAGHU ENGINEERING COLLEGE', font_size=30, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name='Times New Roman')
add_textbox(slide, Inches(1), Inches(1.15), Inches(16), Inches(0.5),
            'Department of Computer Science & Engineering', font_size=20,
            font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name='Times New Roman')

add_textbox(slide, Inches(1.5), Inches(3.3), Inches(15), Inches(1.2),
            'PROJECT MANAGEMENT DASHBOARD', font_size=42,
            font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
add_textbox(slide, Inches(1.5), Inches(4.0), Inches(15), Inches(0.7),
            'Interactive Gantt Chart Visualization System', font_size=24,
            font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_textbox(slide, Inches(1.5), Inches(5.2), Inches(15), Inches(0.5),
            'Tech Stack: React.js  |  Node.js  |  Express.js  |  MongoDB  |  Chart.js  |  DHTMLX Gantt',
            font_size=16, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.2), Inches(8), Inches(1.2),
            'Presented By:\nBatch Members', font_size=16,
            font_color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Times New Roman')
add_textbox(slide, Inches(10), Inches(6.2), Inches(7), Inches(1.2),
            'Under the Guidance of:\nFaculty Guide\nDepartment of CSE', font_size=16,
            font_color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Times New Roman')


# ════════════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'TABLE OF CONTENTS')

contents = [
    '1.  Abstract',
    '2.  Introduction',
    '3.  Existing System & Problems',
    '4.  Proposed System & Advantages',
    '5.  Existing vs Proposed Comparison',
    '6.  System Architecture',
    '7.  Use Case Diagram',
    '8.  Class Diagram',
    '9.  Sequence Diagrams',
    '10. Activity Diagrams',
    '11. Component Diagram',
    '12. ER Diagram',
    '13. State Diagrams',
    '14. Data Flow Diagram',
    '15. Modules & Features',
    '16. API Endpoints',
    '17. Tech Stack',
    '18. Results & Screenshots',
    '19. Future Enhancements',
    '20. Conclusion',
]

# Split into 2 columns
half = len(contents) // 2
add_bullet_list(slide, Inches(1.5), Inches(1.6), Inches(7), Inches(5.5),
                contents[:half], font_size=18, font_color=DARK_GRAY, bullet_char='▸')
add_bullet_list(slide, Inches(9), Inches(1.6), Inches(7), Inches(5.5),
                contents[half:], font_size=18, font_color=DARK_GRAY, bullet_char='▸')


# ════════════════════════════════════════════════════════════════════
# SLIDE 3: ABSTRACT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'ABSTRACT')

abstract_text = (
    "Managing multiple projects, tasks, and resources efficiently is challenging "
    "without proper visualization and tracking tools. Traditional approaches rely on "
    "spreadsheets, emails, and standalone desktop software, leading to data silos, "
    "lack of real-time visibility, and poor collaboration.\n\n"
    "This project presents a comprehensive, full-stack Project Management Dashboard "
    "with Interactive Gantt Chart Visualization, built using the MERN stack "
    "(MongoDB, Express.js, React.js, Node.js). The system provides a centralized "
    "web-based platform for creating and tracking projects, managing tasks with "
    "dependencies and milestones, allocating resources, and visualizing project "
    "timelines through interactive Gantt charts.\n\n"
    "Key features include a real-time analytics dashboard with KPI cards and charts, "
    "drag-and-drop Gantt scheduling, automated project completion calculation, "
    "resource utilization tracking, JWT-based authentication, and a fully responsive "
    "user interface with dark/light theme support."
)
add_textbox(slide, Inches(1), Inches(1.6), Inches(16), Inches(5.5),
            abstract_text, font_size=18, font_color=DARK_GRAY, line_spacing=1.5)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4: INTRODUCTION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'INTRODUCTION')

add_section_title(slide, 'Problem Statement', Inches(1), Inches(1.5))
add_textbox(slide, Inches(1), Inches(2.1), Inches(16), Inches(1.2),
            "Organizations managing multiple projects face challenges in tracking progress, "
            "allocating resources, and meeting deadlines. Without centralized tools, project data "
            "becomes scattered across spreadsheets, emails, and disconnected systems, leading to "
            "poor decision-making and missed deadlines.",
            font_size=16, font_color=DARK_GRAY)

add_section_title(slide, 'Objective', Inches(1), Inches(3.5))
add_bullet_list(slide, Inches(1), Inches(4.1), Inches(16), Inches(3.5), [
    "Build a centralized web-based project management platform",
    "Provide interactive Gantt chart visualization with drag-and-drop scheduling",
    "Enable real-time dashboard analytics with KPIs and charts",
    "Implement task dependency tracking and milestone management",
    "Track resource allocation and workload utilization",
    "Deliver JWT-secured RESTful API architecture",
    "Support responsive design for all device types",
], font_size=16)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5: EXISTING SYSTEM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'EXISTING SYSTEM')

add_section_title(slide, 'Traditional Project Management Methods', Inches(1), Inches(1.5))

table_data = [
    ['Method', 'Description'],
    ['Spreadsheets (Excel)', 'Project timelines & task lists in disconnected files'],
    ['Email-Based Tracking', 'Updates circulated via email threads'],
    ['Desktop Tools (MS Project)', 'Standalone tools without real-time collaboration'],
    ['Physical Whiteboards', 'Sticky notes and printed Gantt charts'],
]
add_table(slide, Inches(1), Inches(2.3), Inches(16), Inches(2.5), table_data,
          col_widths=[Inches(4.5), Inches(11.5)])

add_section_title(slide, 'Problems with Existing System', Inches(1), Inches(5.0))
problems = [
    "No centralized data — scattered across files & emails",
    "No real-time visibility — decisions based on outdated info",
    "Manual progress tracking — error-prone and time-consuming",
    "No task dependency modeling — schedule impacts invisible",
    "Poor resource management — no workload tracking",
    "Static Gantt charts — no drag-and-drop rescheduling",
]
add_bullet_list(slide, Inches(1), Inches(5.6), Inches(16), Inches(2.5),
                problems, font_size=15, font_color=DARK_GRAY, bullet_char='✗')


# ════════════════════════════════════════════════════════════════════
# SLIDE 6: PROPOSED SYSTEM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'PROPOSED SYSTEM')

add_section_title(slide, 'Web-Based Project Management Dashboard', Inches(1), Inches(1.5))
add_textbox(slide, Inches(1), Inches(2.1), Inches(16), Inches(0.9),
            "A full-stack MERN application providing a centralized, real-time platform for "
            "managing projects, tasks, resources, and timelines — replacing manual workflows "
            "with automated, interactive features.",
            font_size=16, font_color=DARK_GRAY)

# Feature cards - 2 rows of 3
cards = [
    ("📊 Dashboard", ["Real-time KPI cards", "Chart.js visualizations", "Activity feed & deadlines"]),
    ("📁 Project Management", ["Full CRUD operations", "Search, filter, sort, paginate", "Auto-delay detection"]),
    ("✅ Task Management", ["Dependencies & milestones", "Resource assignment", "Progress tracking"]),
    ("📅 Gantt Chart", ["Drag-and-drop scheduling", "dhtmlx-gantt library", "Zoom levels & dependency arrows"]),
    ("👥 Resource Management", ["Workload utilization", "Skills & availability", "Auto-computed metrics"]),
    ("🔐 Authentication", ["JWT token-based auth", "bcrypt password hashing", "Role-based access control"]),
]

for i, (title, items) in enumerate(cards):
    col = i % 3
    row = i // 3
    left = Inches(1) + col * Inches(5.3)
    top = Inches(3.2) + row * Inches(2.2)
    add_info_card(slide, left, top, Inches(5), Inches(2.0), title, items)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7: EXISTING VS PROPOSED COMPARISON
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'EXISTING vs PROPOSED SYSTEM')

comp_data = [
    ['Feature', 'Existing System', 'Proposed System'],
    ['Data Storage', 'Scattered files & spreadsheets', 'Centralized MongoDB database'],
    ['Real-Time Updates', '❌ None', '✅ Instant via REST API'],
    ['Gantt Charts', 'Static (Excel)', 'Interactive drag-and-drop'],
    ['Task Dependencies', 'Not tracked', 'Modeled & visualized'],
    ['Resource Tracking', 'Manual calculation', 'Auto-computed utilization'],
    ['Authentication', 'File sharing permissions', 'JWT + role-based control'],
    ['Dashboard / KPIs', 'Manual reports', 'Auto-generated real-time charts'],
    ['Mobile Access', '❌ Limited', '✅ Responsive design'],
    ['Scalability', 'Poor (file-based)', 'High (MERN stack)'],
    ['Notifications', 'Email only', 'In-app toast notifications'],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(16.4), Inches(5.8), comp_data,
          col_widths=[Inches(3.5), Inches(6.3), Inches(6.6)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 8: SYSTEM ARCHITECTURE (HIGH LEVEL)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'SYSTEM ARCHITECTURE — Three-Tier')

# Presentation Tier
add_shape_with_fill(slide, Inches(3), Inches(1.6), Inches(12), Inches(1.5), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(3.3), Inches(1.65), Inches(11.4), Inches(0.4),
            'PRESENTATION TIER', font_size=18, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.3), Inches(2.1), Inches(11.4), Inches(0.9),
            'React.js 18 SPA  •  React Router v6  •  Chart.js  •  dhtmlx-gantt\n'
            'Context API (State)  •  Axios HTTP Client  •  Framer Motion  •  Port 3000',
            font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down
add_textbox(slide, Inches(8.5), Inches(3.15), Inches(1.5), Inches(0.5),
            '⬇  REST API / JSON + JWT  ⬇', font_size=12, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

# Application Tier
add_shape_with_fill(slide, Inches(3), Inches(3.7), Inches(12), Inches(1.5), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(3.3), Inches(3.75), Inches(11.4), Inches(0.4),
            'APPLICATION TIER', font_size=18, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.3), Inches(4.2), Inches(11.4), Inches(0.9),
            'Node.js + Express.js (Port 5000)  •  JWT Auth Middleware  •  CORS  •  Error Handler\n'
            'Controllers: Auth | Project | Task | Resource | Dashboard',
            font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down
add_textbox(slide, Inches(8.5), Inches(5.25), Inches(1.5), Inches(0.5),
            '⬇  Mongoose ODM  ⬇', font_size=12, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

# Data Tier
add_shape_with_fill(slide, Inches(3), Inches(5.8), Inches(12), Inches(1.5), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(3.3), Inches(5.85), Inches(11.4), Inches(0.4),
            'DATA TIER', font_size=18, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.3), Inches(6.3), Inches(11.4), Inches(0.9),
            'MongoDB Atlas / Local Instance  •  In-Memory MongoDB Fallback\n'
            'Collections: Users | Projects | Tasks | Resources',
            font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Labels on left
add_textbox(slide, Inches(0.3), Inches(1.8), Inches(2.5), Inches(0.5),
            '🖥️ Frontend', font_size=16, font_color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(0.3), Inches(3.95), Inches(2.5), Inches(0.5),
            '⚙️ Backend', font_size=16, font_color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(0.3), Inches(6.05), Inches(2.5), Inches(0.5),
            '🗄️ Database', font_size=16, font_color=DARK_BLUE, bold=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9: CLIENT-SIDE ARCHITECTURE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'CLIENT-SIDE ARCHITECTURE')

add_textbox(slide, Inches(1), Inches(1.5), Inches(16), Inches(0.5),
            'React.js Frontend — Component-Based SPA Architecture',
            font_size=20, font_color=DARK_BLUE, bold=True)

arch_text = (
    "index.js  →  App.js (Router + Layout)\n"
    "├── Layout: Navbar.js + Sidebar.js\n"
    "├── Pages: Dashboard | Projects | ProjectDetails | Tasks | GanttChart | Resources | Profile | Settings | Help\n"
    "├── Common: Badge | Loading | Modal | ProgressBar\n"
    "├── State: AppContext.js (React Context API — projects, tasks, resources, dashboard)\n"
    "└── Service: api.js (Axios — request/response interceptors, JWT token injection)"
)
add_textbox(slide, Inches(1.2), Inches(2.2), Inches(15.5), Inches(2.5),
            arch_text, font_size=15, font_color=DARK_GRAY, font_name='Consolas', line_spacing=1.6)

add_section_title(slide, 'Third-Party Libraries', Inches(1), Inches(4.8))
libs_data = [
    ['Library', 'Purpose'],
    ['React 18', 'UI framework with hooks & context'],
    ['React Router v6', 'Client-side navigation & routing'],
    ['Chart.js + react-chartjs-2', 'Doughnut, bar, and line charts'],
    ['dhtmlx-gantt', 'Interactive Gantt chart with drag & drop'],
    ['Axios', 'HTTP client with interceptors'],
    ['Framer Motion', 'Page transitions & animations'],
    ['react-hot-toast', 'Toast notifications'],
    ['date-fns', 'Date formatting & calculations'],
]
add_table(slide, Inches(1), Inches(5.4), Inches(12), Inches(2.3), libs_data,
          col_widths=[Inches(4), Inches(8)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 10: SERVER-SIDE ARCHITECTURE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'SERVER-SIDE ARCHITECTURE')

add_textbox(slide, Inches(1), Inches(1.5), Inches(16), Inches(0.5),
            'Node.js + Express.js — Layered REST API Architecture',
            font_size=20, font_color=DARK_BLUE, bold=True)

server_arch = (
    "server.js (Entry Point — Express Init, Middleware, Route Mounting)\n"
    "│\n"
    "├── Middleware:  CORS  →  JSON Parser  →  Morgan Logger  →  JWT Auth  →  Error Handler\n"
    "│\n"
    "├── Routes:     /api/auth  |  /api/projects  |  /api/tasks  |  /api/resources  |  /api/dashboard\n"
    "│\n"
    "├── Controllers: authController  |  projectController  |  taskController  |  resourceController  |  dashboardController\n"
    "│\n"
    "├── Models:     User  |  Project  |  Task  |  Resource  (Mongoose Schemas + Virtuals + Hooks)\n"
    "│\n"
    "└── Utilities:  asyncHandler  |  ErrorResponse  |  seedData"
)
add_textbox(slide, Inches(1.2), Inches(2.2), Inches(15.5), Inches(3.0),
            server_arch, font_size=14, font_color=DARK_GRAY, font_name='Consolas', line_spacing=1.5)

add_section_title(slide, 'Backend Technologies', Inches(1), Inches(5.3))
be_data = [
    ['Technology', 'Purpose'],
    ['Node.js', 'JavaScript runtime environment'],
    ['Express.js 4', 'Web framework with middleware pipeline'],
    ['MongoDB + Mongoose 8', 'NoSQL database with ODM'],
    ['JWT (jsonwebtoken)', 'Stateless authentication tokens'],
    ['bcryptjs', 'Password hashing (salt rounds: 10)'],
    ['express-validator', 'Request body validation'],
    ['morgan', 'HTTP request logging'],
]
add_table(slide, Inches(1), Inches(5.9), Inches(12), Inches(2.0), be_data,
          col_widths=[Inches(4), Inches(8)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 11: USE CASE DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'USE CASE DIAGRAM')

# Actors column
add_section_title(slide, 'Actors', Inches(0.5), Inches(1.4), Inches(3))
actors = [
    "👤 Project Manager — Full access",
    "👤 Team Member — View + update tasks",
    "👤 Admin — System administration",
]
add_bullet_list(slide, Inches(0.5), Inches(2.0), Inches(3.5), Inches(2.5),
                actors, font_size=13, font_color=DARK_GRAY, bullet_char='')

# Use cases in grouped boxes
uc_groups = [
    ("Authentication", ["Register", "Login", "Logout"], Inches(4.5), Inches(1.4)),
    ("Project Mgmt", ["Create Project", "View Projects", "Update Project", "Delete Project", "View Details", "Search/Filter"], Inches(9), Inches(1.4)),
    ("Task Mgmt", ["Create Task", "Update Task", "Delete Task", "Assign Resource", "Set Dependency", "Mark Milestone", "Track Progress"], Inches(4.5), Inches(3.8)),
    ("Gantt Chart", ["View Gantt", "Drag & Drop", "View Dependencies"], Inches(9), Inches(3.8)),
    ("Resource Mgmt", ["Add Resource", "View Resources", "Update Resource", "View Utilization"], Inches(13.5), Inches(1.4)),
    ("Settings", ["Manage Profile", "Toggle Theme"], Inches(13.5), Inches(3.8)),
]
for title, cases, left, top in uc_groups:
    box_h = Inches(0.3 + len(cases) * 0.3)
    add_shape_with_fill(slide, left, top, Inches(3.8), box_h, BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), Inches(3.6), Inches(0.3),
                title, font_size=14, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, uc in enumerate(cases):
        add_textbox(slide, left + Inches(0.2), top + Inches(0.35 + j * 0.28), Inches(3.4), Inches(0.28),
                    f"○ {uc}", font_size=11, font_color=DARK_GRAY)

add_textbox(slide, Inches(0.5), Inches(7.0), Inches(16), Inches(0.4),
            '28 Use Cases  •  3 Actor Types  •  6 Functional Modules',
            font_size=14, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12: CLASS DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'CLASS DIAGRAM')

classes = [
    ("User", Inches(0.5), Inches(1.5), [
        "_id: ObjectId", "name: String", "email: String",
        "password: String", "role: user|admin",
        "avatar: String", "isActive: Boolean",
    ], ["matchPassword()", "getSignedJwtToken()"]),
    ("Project", Inches(4.8), Inches(1.5), [
        "_id: ObjectId", "name: String", "description: String",
        "startDate: Date", "endDate: Date",
        "status: Enum", "completionPercentage: Number",
        "priority: Enum", "manager: String", "budget: Number",
    ], ["daysRemaining()", "totalDuration()"]),
    ("Task", Inches(9.5), Inches(1.5), [
        "_id: ObjectId", "projectId: ObjectId → Project",
        "title: String", "startDate: Date", "endDate: Date",
        "status: Enum", "dependencyTaskId: ObjectId → Task",
        "assignedResource: ObjectId → Resource",
        "milestone: Boolean", "progress: Number", "priority: Enum",
    ], ["duration()", "isOverdue()"]),
    ("Resource", Inches(14), Inches(1.5), [
        "_id: ObjectId", "name: String", "email: String",
        "role: String", "department: String",
        "availability: Boolean", "skills: String[]",
        "maxHoursPerWeek: Number", "currentWorkload: Number",
    ], ["utilizationPercentage()", "availabilityStatus()", "calculateUtilization()"]),
]

for name, left, top, attrs, methods in classes:
    total_lines = 1 + len(attrs) + 1 + len(methods)
    box_h = Inches(0.25 * total_lines + 0.3)
    add_shape_with_fill(slide, left, top, Inches(4), box_h, WHITE, MSO_SHAPE.RECTANGLE)

    # Class name header
    name_box = add_shape_with_fill(slide, left, top, Inches(4), Inches(0.4), DARK_BLUE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.02), Inches(3.8), Inches(0.35),
                name, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    y_offset = top + Inches(0.45)
    for attr in attrs:
        add_textbox(slide, left + Inches(0.1), y_offset, Inches(3.8), Inches(0.22),
                    f"- {attr}", font_size=9, font_color=DARK_GRAY, font_name='Consolas')
        y_offset += Inches(0.22)

    # Divider line
    add_shape_with_fill(slide, left + Inches(0.1), y_offset, Inches(3.8), Inches(0.02), MEDIUM_BLUE)
    y_offset += Inches(0.08)

    for method in methods:
        add_textbox(slide, left + Inches(0.1), y_offset, Inches(3.8), Inches(0.22),
                    f"+ {method}", font_size=9, font_color=MEDIUM_BLUE, font_name='Consolas')
        y_offset += Inches(0.22)

# Relationships
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(17), Inches(0.5),
            'Relationships:   Project ──[1:*]──► Task      Task ──[*:0..1]──► Task (dependency)      Task ──[*:0..1]──► Resource (assignment)',
            font_size=13, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13: SEQUENCE DIAGRAM — LOGIN
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'SEQUENCE DIAGRAM — User Login')

components = ['User', 'React App', 'Axios Client', 'Express Server', 'AuthController', 'MongoDB']
x_positions = [Inches(0.5), Inches(3), Inches(5.5), Inches(8.5), Inches(11.5), Inches(15)]

for comp, xp in zip(components, x_positions):
    add_shape_with_fill(slide, xp, Inches(1.5), Inches(2.2), Inches(0.45), DARK_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, xp, Inches(1.52), Inches(2.2), Inches(0.4),
                comp, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

steps = [
    (0, 1, '1. Enter email & password'),
    (1, 2, '2. POST /api/auth/login'),
    (2, 3, '3. HTTP Request + JSON'),
    (3, 4, '4. Route to login()'),
    (4, 5, '5. User.findOne({email})'),
    (5, 4, '6. User document'),
    (4, 4, '7. matchPassword() → getSignedJwtToken()'),
    (4, 3, '8. {success, token: JWT}'),
    (3, 2, '9. 200 OK + JSON'),
    (2, 1, '10. Store token in localStorage'),
    (1, 0, '11. Redirect to Dashboard'),
]

y = Inches(2.2)
for frm, to, label in steps:
    y += Inches(0.4)
    fx = x_positions[frm] + Inches(1.1)
    tx = x_positions[to] + Inches(1.1)
    direction = '→' if to > frm else '←' if to < frm else '↺'
    add_textbox(slide, min(fx, tx), y, abs(fx - tx) if fx != tx else Inches(2.2), Inches(0.3),
                f"{direction} {label}", font_size=9, font_color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════
# SLIDE 14: SEQUENCE DIAGRAM — CREATE PROJECT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'SEQUENCE DIAGRAM — Create Project')

components2 = ['Manager', 'React App', 'AppContext', 'Axios API', 'Express Server', 'ProjectController', 'MongoDB']
x_pos2 = [Inches(0.3), Inches(2.3), Inches(4.5), Inches(6.7), Inches(9), Inches(11.5), Inches(14.5)]

for comp, xp in zip(components2, x_pos2):
    add_shape_with_fill(slide, xp, Inches(1.5), Inches(2), Inches(0.45), DARK_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, xp, Inches(1.52), Inches(2), Inches(0.4),
                comp, font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

steps2 = [
    (0, 1, '1. Fill project form & submit'),
    (1, 2, '2. createProject(data)'),
    (2, 3, '3. projectAPI.create(data)'),
    (3, 4, '4. POST /api/projects + JSON'),
    (4, 5, '5. Route to createProject()'),
    (5, 6, '6. new Project(data) → save()'),
    (6, 5, '7. Saved document'),
    (5, 4, '8. {success: true, data: project}'),
    (4, 3, '9. 201 Created'),
    (3, 2, '10. Response data'),
    (2, 1, '11. setProjects([new, ...prev])'),
    (1, 0, '12. Toast: "Project created"'),
]

y = Inches(2.2)
for frm, to, label in steps2:
    y += Inches(0.38)
    fx = x_pos2[frm] + Inches(1)
    tx = x_pos2[to] + Inches(1)
    direction = '→' if to > frm else '←'
    add_textbox(slide, min(fx, tx), y, abs(fx - tx) if fx != tx else Inches(2), Inches(0.28),
                f"{direction} {label}", font_size=9, font_color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════
# SLIDE 15: ACTIVITY DIAGRAM — PROJECT LIFECYCLE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'ACTIVITY DIAGRAM — Project Lifecycle')

act_steps = [
    ('●', 'Start', DARK_BLUE),
    ('▼', 'Manager Creates Project', MEDIUM_BLUE),
    ('◇', 'Valid Data?', ORANGE),
    ('▼', 'Save to DB (Status: Not Started)', ACCENT_GREEN),
    ('▼', 'Add Tasks to Project', MEDIUM_BLUE),
    ('▼', 'Assign Resources to Tasks', MEDIUM_BLUE),
    ('▼', 'Project Status: In Progress', LIGHT_BLUE),
    ('◇', 'All Tasks Completed?', ORANGE),
    ('◇', 'Past End Date?', ORANGE),
    ('▼', 'Status: Completed (100%)', ACCENT_GREEN),
    ('▼', 'Status: Delayed → Reschedule', ACCENT_RED),
    ('●', 'End', DARK_BLUE),
]

# Layout in a flow - 2 columns
col1_items = act_steps[:7]
col2_items = act_steps[7:]

for i, (sym, text, color) in enumerate(col1_items):
    top = Inches(1.5) + i * Inches(0.8)
    shape = add_shape_with_fill(slide, Inches(2), top, Inches(5), Inches(0.55), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(2.2), top + Inches(0.08), Inches(4.6), Inches(0.4),
                f"{sym}  {text}", font_size=14, font_color=WHITE, bold=True)

for i, (sym, text, color) in enumerate(col2_items):
    top = Inches(1.5) + i * Inches(0.8)
    shape = add_shape_with_fill(slide, Inches(10), top, Inches(6), Inches(0.55), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(10.2), top + Inches(0.08), Inches(5.6), Inches(0.4),
                f"{sym}  {text}", font_size=14, font_color=WHITE, bold=True)

# Flow explanation
add_textbox(slide, Inches(1), Inches(6.8), Inches(16), Inches(0.7),
            'Flow: Create → Add Tasks → Assign Resources → Track Progress → Complete or Reschedule if Delayed\n'
            'Auto-detection: Project marked "Delayed" if endDate < today and status ≠ Completed',
            font_size=13, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16: ACTIVITY DIAGRAM — AUTHENTICATION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'ACTIVITY DIAGRAM — Authentication Flow')

auth_steps_left = [
    ('● Start', DARK_BLUE),
    ('User Opens Application', MEDIUM_BLUE),
    ('◇ Token in LocalStorage?', ORANGE),
    ('Attach Token to Request Header', LIGHT_BLUE),
    ('◇ Token Valid?', ORANGE),
    ('✓ Access Granted — Load Dashboard', ACCENT_GREEN),
]

auth_steps_right = [
    ('No Token / Invalid Token', ACCENT_RED),
    ('Redirect to Login', MEDIUM_BLUE),
    ('Enter Credentials', LIGHT_BLUE),
    ('POST /api/auth/login', MEDIUM_BLUE),
    ('◇ Valid Credentials?', ORANGE),
    ('Generate JWT Token → Store in localStorage', ACCENT_GREEN),
]

for i, (text, color) in enumerate(auth_steps_left):
    top = Inches(1.5) + i * Inches(0.85)
    add_shape_with_fill(slide, Inches(1), top, Inches(6.5), Inches(0.55), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(1.2), top + Inches(0.08), Inches(6.1), Inches(0.4),
                text, font_size=13, font_color=WHITE, bold=True)

for i, (text, color) in enumerate(auth_steps_right):
    top = Inches(1.5) + i * Inches(0.85)
    add_shape_with_fill(slide, Inches(9.5), top, Inches(7.5), Inches(0.55), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(9.7), top + Inches(0.08), Inches(7.1), Inches(0.4),
                text, font_size=13, font_color=WHITE, bold=True)

add_textbox(slide, Inches(7.7), Inches(2.4), Inches(1.6), Inches(0.4),
            'YES →', font_size=12, font_color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.7), Inches(3.2), Inches(1.6), Inches(0.4),
            'NO →', font_size=12, font_color=ACCENT_RED, bold=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 17: COMPONENT DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'COMPONENT DIAGRAM')

# Frontend box
add_shape_with_fill(slide, Inches(0.5), Inches(1.4), Inches(8.5), Inches(5.8), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(8), Inches(0.4),
            '◀ FRONTEND (React.js) ▶', font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

fe_components = [
    ("UI Layer", "Pages: Dashboard, Projects, Tasks, GanttChart,\nResources, Profile, Settings, Help"),
    ("Layout", "Navbar.js + Sidebar.js"),
    ("Common", "Badge, Loading, Modal, ProgressBar"),
    ("State Mgmt", "AppContext.js (React Context API)"),
    ("Service", "api.js (Axios HTTP Client)"),
    ("Libraries", "Chart.js, dhtmlx-gantt, Framer Motion,\nreact-hot-toast, React Router v6, date-fns"),
]
for i, (label, desc) in enumerate(fe_components):
    y = Inches(2.0) + i * Inches(0.82)
    add_textbox(slide, Inches(0.8), y, Inches(2), Inches(0.35),
                label, font_size=12, font_color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(2.8), y, Inches(6), Inches(0.7),
                desc, font_size=11, font_color=DARK_GRAY)

# Backend box
add_shape_with_fill(slide, Inches(9.5), Inches(1.4), Inches(8), Inches(5.8), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(9.7), Inches(1.5), Inches(7.5), Inches(0.4),
            '◀ BACKEND (Express.js) ▶', font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

be_components = [
    ("API Routes", "/auth, /projects, /tasks,\n/resources, /dashboard"),
    ("Middleware", "JWT Auth, Error Handler, CORS"),
    ("Controllers", "Auth, Project, Task,\nResource, Dashboard"),
    ("Models", "User, Project, Task, Resource\n(Mongoose Schemas)"),
    ("Utilities", "asyncHandler, ErrorResponse, seedData"),
    ("Database", "MongoDB Atlas / In-Memory Fallback"),
]
for i, (label, desc) in enumerate(be_components):
    y = Inches(2.0) + i * Inches(0.82)
    add_textbox(slide, Inches(9.8), y, Inches(2), Inches(0.35),
                label, font_size=12, font_color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(11.8), y, Inches(5.5), Inches(0.7),
                desc, font_size=11, font_color=DARK_GRAY)

# Arrow between
add_textbox(slide, Inches(8.5), Inches(4.0), Inches(1.2), Inches(0.5),
            'REST\nAPI', font_size=12, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 18: ER DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'ENTITY-RELATIONSHIP (ER) DIAGRAM')

entities = [
    ("USER", Inches(0.5), Inches(1.5), Inches(4), [
        "PK _id: ObjectId", "name: String", "email: String (unique)",
        "password: String", "role: user | admin",
        "isActive: Boolean", "lastLogin: Date"
    ]),
    ("PROJECT", Inches(5), Inches(1.5), Inches(4), [
        "PK _id: ObjectId", "name: String", "description: String",
        "startDate: Date", "endDate: Date",
        "status: Enum(4)", "completionPercentage: Number",
        "priority: Enum(4)", "manager: String", "budget: Number"
    ]),
    ("TASK", Inches(9.5), Inches(1.5), Inches(4.3), [
        "PK _id: ObjectId", "FK projectId → PROJECT",
        "title: String", "startDate: Date", "endDate: Date",
        "status: Enum(4)", "FK dependencyTaskId → TASK",
        "FK assignedResource → RESOURCE",
        "milestone: Boolean", "progress: Number", "priority: Enum(4)"
    ]),
    ("RESOURCE", Inches(14.2), Inches(1.5), Inches(3.8), [
        "PK _id: ObjectId", "name: String", "email: String",
        "role: String", "department: String",
        "availability: Boolean", "skills: String[]",
        "maxHoursPerWeek: Number", "currentWorkload: Number"
    ]),
]

for name, left, top, width, fields in entities:
    h = Inches(0.28 * len(fields) + 0.6)
    add_shape_with_fill(slide, left, top, width, h, WHITE, MSO_SHAPE.RECTANGLE)
    name_bar = add_shape_with_fill(slide, left, top, width, Inches(0.4), DARK_BLUE)
    add_textbox(slide, left, top + Inches(0.02), width, Inches(0.35),
                name, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, field in enumerate(fields):
        color = MEDIUM_BLUE if field.startswith('PK') or field.startswith('FK') else DARK_GRAY
        add_textbox(slide, left + Inches(0.1), top + Inches(0.45 + j * 0.26),
                    width - Inches(0.2), Inches(0.24),
                    field, font_size=10, font_color=color, font_name='Consolas')

# Relationships text
add_textbox(slide, Inches(0.5), Inches(6.2), Inches(17), Inches(1.0),
            'Relationships:\n'
            'PROJECT ──── [1 : Many] ────► TASK  (A project has many tasks)\n'
            'TASK ──── [Many : 0..1] ────► TASK  (A task may depend on another task)\n'
            'TASK ──── [Many : 0..1] ────► RESOURCE  (A task may be assigned to a resource)',
            font_size=13, font_color=DARK_BLUE, font_name='Calibri', line_spacing=1.5)


# ════════════════════════════════════════════════════════════════════
# SLIDE 19: STATE DIAGRAMS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'STATE DIAGRAMS')

# Project States
add_section_title(slide, 'Project State Machine', Inches(0.5), Inches(1.4), Inches(8))

proj_states = [
    ('Not Started', 'completionPercentage = 0', MEDIUM_GRAY, Inches(1), Inches(2.2)),
    ('In Progress', '0 < completion < 100', MEDIUM_BLUE, Inches(1), Inches(3.5)),
    ('Delayed', 'endDate < today', ACCENT_RED, Inches(4.5), Inches(4.8)),
    ('Completed', 'completion = 100%', ACCENT_GREEN, Inches(1), Inches(4.8)),
]

for name, desc, color, left, top in proj_states:
    add_shape_with_fill(slide, left, top, Inches(3.2), Inches(0.8), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), Inches(3), Inches(0.35),
                name, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.4), Inches(3), Inches(0.3),
                desc, font_size=10, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# Transitions text
add_textbox(slide, Inches(0.5), Inches(5.9), Inches(8), Inches(1.3),
            'Transitions:\n'
            '• Not Started → In Progress (first task started)\n'
            '• In Progress → Completed (all tasks done)\n'
            '• In Progress → Delayed (past end date)\n'
            '• Delayed → In Progress (rescheduled)',
            font_size=11, font_color=DARK_GRAY)

# Task States
add_section_title(slide, 'Task State Machine', Inches(9), Inches(1.4), Inches(8))

task_states = [
    ('Not Started', 'progress = 0', MEDIUM_GRAY, Inches(9.5), Inches(2.2)),
    ('In Progress', '0 < progress < 100', MEDIUM_BLUE, Inches(9.5), Inches(3.5)),
    ('Delayed', 'endDate < today', ACCENT_RED, Inches(13), Inches(4.8)),
    ('Completed', 'progress = 100%', ACCENT_GREEN, Inches(9.5), Inches(4.8)),
]

for name, desc, color, left, top in task_states:
    add_shape_with_fill(slide, left, top, Inches(3.2), Inches(0.8), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), Inches(3), Inches(0.35),
                name, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.4), Inches(3), Inches(0.3),
                desc, font_size=10, font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(9), Inches(5.9), Inches(8), Inches(1.3),
            'Transitions:\n'
            '• Not Started → In Progress (work begins)\n'
            '• In Progress → Completed (progress = 100%)\n'
            '• In Progress → Delayed (past end date)\n'
            '• Delayed → Completed (completed late)',
            font_size=11, font_color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 20: DATA FLOW DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'DATA FLOW DIAGRAM')

# Level 0 - Context
add_section_title(slide, 'Level 0 — Context Diagram', Inches(0.5), Inches(1.4), Inches(8))

# External entity
add_shape_with_fill(slide, Inches(0.5), Inches(2.5), Inches(3), Inches(1.5), MEDIUM_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(0.7), Inches(2.6), Inches(2.6), Inches(1.2),
            '👤 User\n\nProject Data\nTask Data\nCredentials',
            font_size=12, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# System
add_shape_with_fill(slide, Inches(5), Inches(2.2), Inches(5), Inches(2.2), DARK_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(5.2), Inches(2.3), Inches(4.6), Inches(2.0),
            'Project Management\nDashboard System\n\n• Manage Projects\n• Manage Tasks\n• Manage Resources\n• Generate Dashboard',
            font_size=12, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# Outputs
add_shape_with_fill(slide, Inches(11.5), Inches(2.5), Inches(3), Inches(1.5), ACCENT_GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(11.7), Inches(2.6), Inches(2.6), Inches(1.2),
            '📊 Outputs\n\nDashboard KPIs\nGantt Charts\nReports',
            font_size=12, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# Arrows
add_textbox(slide, Inches(3.5), Inches(2.8), Inches(1.5), Inches(0.4),
            '→→→', font_size=18, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(10), Inches(2.8), Inches(1.5), Inches(0.4),
            '→→→', font_size=18, font_color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Level 1
add_section_title(slide, 'Level 1 — Process Diagram', Inches(0.5), Inches(5.0), Inches(8))

processes = [
    ('1.0 Authenticate', Inches(0.5)),
    ('2.0 Manage Projects', Inches(4)),
    ('3.0 Manage Tasks', Inches(7.5)),
    ('4.0 Manage Resources', Inches(11)),
    ('5.0 Generate Dashboard', Inches(14.5)),
]
for pname, left in processes:
    add_shape_with_fill(slide, left, Inches(5.8), Inches(3), Inches(0.7), MEDIUM_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.1), Inches(5.85), Inches(2.8), Inches(0.6),
                pname, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Data store
add_shape_with_fill(slide, Inches(3), Inches(6.8), Inches(12), Inches(0.6), DARK_BLUE, MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(3.2), Inches(6.83), Inches(11.6), Inches(0.5),
            'D1: MongoDB Database  [Users | Projects | Tasks | Resources]',
            font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 21: MODULES & FEATURES
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'MODULES & KEY FEATURES')

modules = [
    ("📊 Dashboard", ["Real-time KPI cards (6 metrics)", "Task status doughnut chart", "Project progress bar chart", "Recent activities & deadlines", "Monthly trend line chart"]),
    ("📁 Projects", ["Create/Edit/Delete projects", "Search, filter, sort, paginate", "Priority & status badges", "Budget & timeline tracking", "Auto-delay detection"]),
    ("✅ Tasks", ["CRUD per project + global view", "Dependency & milestone flags", "Resource assignment", "Progress % & hours tracking", "Reorder tasks"]),
    ("📅 Gantt Chart", ["dhtmlx-gantt library", "Drag-and-drop scheduling", "Dependency arrows", "Milestone diamonds", "Zoom: Hours/Days/Weeks/Months"]),
    ("👥 Resources", ["Team member profiles", "Skills & department info", "Utilization % auto-computed", "Availability status", "Hourly rate & workload"]),
    ("🔐 Auth & Settings", ["JWT token authentication", "bcrypt password hashing", "Role-based access (user/admin)", "Profile management", "Dark/Light theme toggle"]),
]

for i, (title, features) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(5.8)
    top = Inches(1.5) + row * Inches(3.0)
    add_info_card(slide, left, top, Inches(5.5), Inches(2.7), title, features, DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 22: API ENDPOINTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'API ENDPOINTS')

api_data = [
    ['Method', 'Endpoint', 'Description'],
    ['POST', '/api/auth/register', 'Register new user'],
    ['POST', '/api/auth/login', 'Login & get JWT token'],
    ['GET', '/api/auth/me', 'Get current user'],
    ['GET', '/api/projects', 'List all projects'],
    ['POST', '/api/projects', 'Create project'],
    ['PUT', '/api/projects/:id', 'Update project'],
    ['DELETE', '/api/projects/:id', 'Delete project'],
    ['GET', '/api/tasks/:projectId', 'Get project tasks'],
    ['GET', '/api/tasks/all', 'Get all tasks'],
    ['POST', '/api/tasks', 'Create task'],
    ['PUT', '/api/tasks/:id', 'Update task'],
    ['DELETE', '/api/tasks/:id', 'Delete task'],
    ['GET', '/api/tasks/gantt/:id', 'Get Gantt chart data'],
    ['GET', '/api/resources', 'List all resources'],
    ['POST', '/api/resources', 'Add resource'],
    ['PUT', '/api/resources/:id', 'Update resource'],
    ['DELETE', '/api/resources/:id', 'Delete resource'],
    ['GET', '/api/dashboard/overview', 'Dashboard statistics'],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(16.4), Inches(5.8), api_data,
          col_widths=[Inches(2), Inches(6), Inches(8.4)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 23: TECH STACK SUMMARY
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'TECHNOLOGY STACK')

# Frontend
add_section_title(slide, 'Frontend', Inches(0.5), Inches(1.5), Inches(8))
fe_stack = [
    ['Technology', 'Version', 'Purpose'],
    ['React.js', '18.x', 'UI library with hooks & context'],
    ['React Router', 'v6', 'Client-side routing'],
    ['Chart.js', '4.x', 'Doughnut, bar & line charts'],
    ['dhtmlx-gantt', 'Latest', 'Interactive Gantt chart'],
    ['Axios', '1.x', 'HTTP client with interceptors'],
    ['Framer Motion', 'Latest', 'Page animations'],
    ['react-hot-toast', 'Latest', 'Toast notifications'],
    ['date-fns', 'Latest', 'Date formatting'],
]
add_table(slide, Inches(0.5), Inches(2.1), Inches(8), Inches(3.0), fe_stack,
          col_widths=[Inches(2.5), Inches(1.5), Inches(4)])

# Backend
add_section_title(slide, 'Backend', Inches(9.5), Inches(1.5), Inches(8))
be_stack = [
    ['Technology', 'Version', 'Purpose'],
    ['Node.js', '≥16', 'JavaScript runtime'],
    ['Express.js', '4.x', 'Web framework'],
    ['MongoDB', 'Latest', 'NoSQL database'],
    ['Mongoose', '8.x', 'ODM for MongoDB'],
    ['JWT', '9.x', 'Authentication tokens'],
    ['bcryptjs', '2.x', 'Password hashing'],
    ['morgan', '1.x', 'HTTP logging'],
    ['mongodb-memory-server', '9.x', 'In-memory DB fallback'],
]
add_table(slide, Inches(9.5), Inches(2.1), Inches(8), Inches(3.0), be_stack,
          col_widths=[Inches(3), Inches(1.5), Inches(3.5)])

# Dev tools
add_section_title(slide, 'Dev Tools', Inches(0.5), Inches(5.5), Inches(8))
dev_items = [
    "concurrently — Run frontend + backend simultaneously",
    "nodemon — Auto-restart server on code changes",
    "VS Code — IDE with React & Node.js extensions",
    "Postman — API testing",
    "Git — Version control",
]
add_bullet_list(slide, Inches(0.5), Inches(6.1), Inches(16), Inches(1.5),
                dev_items, font_size=14, font_color=DARK_GRAY, bullet_char='▸')


# ════════════════════════════════════════════════════════════════════
# SLIDE 24: RESULTS & SCREENSHOTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'RESULTS & SCREENSHOTS')

screenshots = [
    ("Dashboard View", "Real-time KPI cards, task distribution chart,\nproject progress bars, recent activities"),
    ("Projects Page", "Project cards with status badges,\nsearch & filter controls, CRUD operations"),
    ("Gantt Chart", "Interactive timeline with drag-and-drop,\ntask dependencies and milestones"),
    ("Tasks Page", "Task list with filters, progress tracking,\nassigned resources and priority badges"),
    ("Resources Page", "Team member cards, utilization bars,\navailability status and skills display"),
    ("Project Details", "Full project info, task breakdown,\ncompletion metrics and quick actions"),
]

for i, (title, desc) in enumerate(screenshots):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(5.8)
    top = Inches(1.5) + row * Inches(3.0)

    # Placeholder box for screenshot
    add_shape_with_fill(slide, left, top, Inches(5.5), Inches(1.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.5), top + Inches(0.4), Inches(4.5), Inches(1.0),
                '[ Screenshot Placeholder ]', font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, left, top + Inches(1.9), Inches(5.5), Inches(0.35),
                title, font_size=15, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(2.25), Inches(5.5), Inches(0.7),
                desc, font_size=11, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 25: FUTURE ENHANCEMENTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'FUTURE ENHANCEMENTS')

enhancements = [
    ("📧 Email Notifications", ["Deadline reminders", "Task assignment alerts", "Weekly digest reports"]),
    ("📊 Advanced Reporting", ["PDF/Excel export", "Custom report builder", "Budget vs actual analysis"]),
    ("💬 Collaboration", ["Comments on tasks", "File attachments", "Real-time team chat"]),
    ("📱 Mobile Application", ["React Native version", "Offline support", "Push notifications"]),
    ("🤖 AI Features", ["Automated task estimation", "Risk prediction", "Resource optimization"]),
    ("📈 Advanced Analytics", ["Burndown charts", "Velocity tracking", "Predictive timelines"]),
]

for i, (title, items) in enumerate(enhancements):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(5.8)
    top = Inches(1.5) + row * Inches(3.0)
    add_info_card(slide, left, top, Inches(5.5), Inches(2.5), title, items, DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 26: CONCLUSION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, 'CONCLUSION')

add_section_title(slide, 'Project Achievements', Inches(1), Inches(1.5))
achievements = [
    "Fully functional project management system with MERN stack",
    "Interactive Gantt chart visualization with drag-and-drop",
    "Real-time dashboard analytics with KPI cards and charts",
    "Complete CRUD operations for projects, tasks, and resources",
    "JWT-based authentication with role-based access control",
    "Responsive design accessible on all devices",
    "Clean, modern UI with dark/light theme support",
]
add_bullet_list(slide, Inches(1), Inches(2.1), Inches(7.5), Inches(3.0),
                achievements, font_size=15, font_color=DARK_GRAY, bullet_char='✓')

add_section_title(slide, 'Learning Outcomes', Inches(9.5), Inches(1.5))
learnings = [
    "Full-stack development with MERN stack",
    "RESTful API design and implementation",
    "Chart.js and Gantt chart library integration",
    "React Context API for state management",
    "MongoDB data modeling with Mongoose",
    "JWT authentication implementation",
    "Responsive CSS and component design",
]
add_bullet_list(slide, Inches(9.5), Inches(2.1), Inches(7.5), Inches(3.0),
                learnings, font_size=15, font_color=DARK_GRAY, bullet_char='▸')

add_section_title(slide, 'Business Value', Inches(1), Inches(5.3))
biz_value = [
    "Improved project visibility and decision-making",
    "Better resource allocation and workload balancing",
    "Enhanced team collaboration through centralized platform",
    "Data-driven insights through automated analytics",
]
add_bullet_list(slide, Inches(1), Inches(5.9), Inches(16), Inches(2.0),
                biz_value, font_size=15, font_color=DARK_GRAY, bullet_char='★')


# ════════════════════════════════════════════════════════════════════
# SLIDE 27: THANK YOU
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_shape_with_fill(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, DARK_BLUE)
add_shape_with_fill(slide, Emu(0), Inches(3.2), SLIDE_W, Inches(0.06), LIGHT_BLUE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(16), Inches(1.2),
            'THANK YOU', font_size=60, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_textbox(slide, Inches(1), Inches(3.5), Inches(16), Inches(0.8),
            'Project Management Dashboard\nInteractive Gantt Chart Visualization System',
            font_size=24, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.0), Inches(16), Inches(0.5),
            'Tech Stack: React.js  |  Node.js  |  Express.js  |  MongoDB  |  Chart.js  |  DHTMLX Gantt',
            font_size=16, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(5), Inches(6.5), Inches(8), Inches(0.8),
            'Q & A',
            font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"✅ Presentation saved to: {OUTPUT_PATH}")
print(f"📊 Total slides: {len(prs.slides)}")
