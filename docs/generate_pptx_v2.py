"""
Generate Project Management Dashboard PPTX
Styled to match sampleppt (2).pptx — clean Canva-like white slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Use sample (2) as dimensional reference ───
SAMPLE_PATH = r"c:\Users\NSVLJ\Downloads\sampleppt (2).pptx"
OUTPUT_PATH = r"c:\Projects\project-management-dashboard\docs\Project_Management_Dashboard_Presentation.pptx"

sample_prs = Presentation(SAMPLE_PATH)
prs = Presentation()
prs.slide_width = sample_prs.slide_width    # 20.0in
prs.slide_height = sample_prs.slide_height  # 11.2in

blank_layout = prs.slide_layouts[6]  # Blank

# ─── Color Palette (matching sample) ───
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT     = RGBColor(0x25, 0x1C, 0x23)
DARK_BLUE    = RGBColor(0x1E, 0x3A, 0x5F)
MEDIUM_BLUE  = RGBColor(0x3C, 0x6C, 0xA8)
LIGHT_BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY  = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
BG_BLUE      = RGBColor(0xE8, 0xF0, 0xFE)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Title font size matching sample: ~48pt (609600 EMU)
TITLE_SIZE = 48
BODY_SIZE = 32
SMALL_SIZE = 28
TABLE_HEADER_SIZE = 20
TABLE_BODY_SIZE = 18
BULLET_SIZE = 28


def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=BODY_SIZE,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Times New Roman", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_slide_title(slide, text, top=Inches(0.7)):
    """Add centered title matching sample style — large Times New Roman"""
    add_textbox(slide, Inches(2), top, Inches(16), Inches(0.9),
                text, font_size=TITLE_SIZE, font_color=BLACK, bold=False,
                alignment=PP_ALIGN.CENTER, font_name="Times New Roman")


def add_body_text(slide, text, left=Inches(1), top=Inches(2.2), width=Inches(17.5),
                  height=Inches(7), font_size=BODY_SIZE, font_color=BLACK, bold=False,
                  line_spacing=1.5):
    return add_textbox(slide, left, top, width, height, text,
                       font_size=font_size, font_color=font_color, bold=bold,
                       line_spacing=line_spacing)


def add_bullet_textbox(slide, left, top, width, height, items, font_size=BULLET_SIZE,
                       font_color=DARK_TEXT, font_name="Times New Roman", bullet="•",
                       bold_prefix=False):
    """Add a multi-paragraph bulleted textbox — matches sample body style"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(4)

        if bold_prefix and ":" in item:
            # Bold the part before the colon
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = f"{bullet} {prefix}:"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = font_color
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = font_color
            run2.font.bold = False
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = f"{bullet} {item}" if bullet else item
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = False
            run.font.name = font_name
    return txBox


def add_table(slide, left, top, width, height, rows_data, col_widths=None):
    """Add a table with dark blue header"""
    nrows = len(rows_data)
    ncols = len(rows_data[0]) if rows_data else 1
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(TABLE_HEADER_SIZE if ri == 0 else TABLE_BODY_SIZE)
                para.font.name = "Times New Roman"
                if ri == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = DARK_TEXT
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return tbl_shape


def add_decorative_circle(slide, left, top, size, color):
    """Add a decorative circle like the sample uses"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Make semi-transparent via adjusting alpha isn't easy in python-pptx,
    # so we just use a lighter color
    return shape


# ════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Decorative circles (like sample)
add_decorative_circle(slide, Inches(16.5), Inches(-1), Inches(3), LIGHT_BLUE)
add_decorative_circle(slide, Inches(-0.5), Inches(8.5), Inches(2.5), BG_BLUE)

add_textbox(slide, Inches(2), Inches(1.5), Inches(16), Inches(0.8),
            "RAGHU ENGINEERING COLLEGE", font_size=37, font_color=BLACK,
            bold=False, alignment=PP_ALIGN.CENTER, font_name="Times New Roman")

add_textbox(slide, Inches(1.5), Inches(2.5), Inches(17), Inches(0.6),
            '"COMPUTER SCIENCE AND ENGINEERING"', font_size=34,
            font_color=BLACK, alignment=PP_ALIGN.CENTER, font_name="Times New Roman")

add_textbox(slide, Inches(1), Inches(4.2), Inches(18), Inches(1.8),
            "Project Management Dashboard\nwith Interactive Gantt Chart Visualization System",
            font_size=44, font_color=BLACK, bold=False,
            alignment=PP_ALIGN.CENTER, font_name="Times New Roman", line_spacing=1.3)

add_textbox(slide, Inches(4), Inches(6.5), Inches(12), Inches(0.5),
            "Tech Stack: React.js | Node.js | Express.js | MongoDB | Chart.js | DHTMLX Gantt",
            font_size=20, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

# Batch info (left)
add_textbox(slide, Inches(1), Inches(8.0), Inches(6), Inches(2),
            "BATCH - XX\nStudent 1 - Roll No\nStudent 2 - Roll No\nStudent 3 - Roll No",
            font_size=20, font_color=BLACK, font_name="Times New Roman")

# Guide info (right)
add_textbox(slide, Inches(13), Inches(8.0), Inches(6), Inches(2),
            "Under the guidance\nMr./Ms. Faculty Name\nCSE DEPARTMENT",
            font_size=20, font_color=BLACK, font_name="Times New Roman")

# Presented by label (left bottom)
add_textbox(slide, Inches(1), Inches(7.2), Inches(6), Inches(0.5),
            "PRESENTED BY", font_size=22, font_color=MEDIUM_BLUE,
            bold=False, font_name="Times New Roman")


# ════════════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

add_textbox(slide, Inches(1), Inches(1.0), Inches(16), Inches(0.9),
            "TABLE OF CONTENTS", font_size=TITLE_SIZE, font_color=BLACK,
            bold=False, alignment=PP_ALIGN.CENTER, font_name="Times New Roman")

contents = [
    "Abstract",
    "Introduction",
    "Problem Statement",
    "Existing System",
    "Proposed System",
    "Proposed Vs Existing System",
    "System Architecture",
    "UML Diagrams",
    "Modules & Key Features",
    "API Endpoints",
    "Technology Stack",
    "Authentication System",
    "Results & Screenshots",
    "Future Enhancements",
    "Conclusion",
]

add_bullet_textbox(slide, Inches(3), Inches(2.2), Inches(14), Inches(7.5),
                   contents, font_size=SMALL_SIZE, bullet="•")


# ════════════════════════════════════════════════════════════════════
# SLIDE 3: ABSTRACT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "ABSTRACT")

abstract = (
    "Managing multiple projects, tasks, and resources efficiently is challenging "
    "without proper visualization and tracking tools. Traditional approaches rely on "
    "spreadsheets, emails, and standalone desktop software, leading to data silos, "
    "lack of real-time visibility, and poor collaboration.\n\n"
    "This project presents a comprehensive, full-stack Project Management Dashboard "
    "with Interactive Gantt Chart Visualization, built using the MERN stack "
    "(MongoDB, Express.js, React.js, Node.js). The system provides a centralized "
    "web-based platform for creating and tracking projects, managing tasks with "
    "dependencies and milestones, allocating resources, and visualizing project "
    "timelines through interactive Gantt charts.\n\n"
    "Key features include a real-time analytics dashboard with KPI cards and charts, "
    "drag-and-drop Gantt scheduling, automated project completion calculation, "
    "resource utilization tracking, JWT-based authentication with OTP password reset, "
    "email notifications via Nodemailer, and a fully responsive user interface."
)
add_body_text(slide, abstract, font_size=BODY_SIZE, line_spacing=1.4)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4: INTRODUCTION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "INTRODUCTION")

intro = (
    "In today's fast-paced business environment, organizations manage multiple projects simultaneously "
    "across diverse teams and departments. The ability to track progress, allocate resources efficiently, "
    "and meet deadlines is critical for organizational success.\n\n"
    "Without centralized tools, project data becomes scattered across spreadsheets, emails, and "
    "disconnected systems, leading to poor decision-making and missed deadlines. Project managers "
    "need a single platform that provides real-time visibility into project health, task status, "
    "resource utilization, and timeline visualization.\n\n"
    "This project addresses these challenges by building a modern, web-based Project Management "
    "Dashboard using the MERN stack (MongoDB, Express.js, React.js, Node.js) with interactive "
    "Gantt chart visualization, comprehensive authentication, and real-time analytics."
)
add_body_text(slide, intro, font_size=BODY_SIZE, line_spacing=1.4)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "PROBLEM STATEMENT")

problem = (
    "Organizations managing multiple projects face several critical challenges:\n\n"
    "• No centralized data — project information is scattered across files, spreadsheets, and emails\n\n"
    "• No real-time visibility — decisions are based on outdated information from manual reports\n\n"
    "• Manual progress tracking — error-prone and time-consuming status updates\n\n"
    "• No task dependency modeling — schedule impacts from delays are invisible\n\n"
    "• Poor resource management — no automated workload tracking or utilization metrics\n\n"
    "• Static Gantt charts — Excel-based timelines with no drag-and-drop rescheduling\n\n"
    "• No authentication or access control — anyone with the file can modify project data"
)
add_body_text(slide, problem, font_size=SMALL_SIZE, line_spacing=1.35)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6: EXISTING SYSTEM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Existing System")

existing_items = [
    "Spreadsheets (Excel): Project timelines and task lists in disconnected files with no real-time collaboration",
    "Email-Based Tracking: Project updates circulated via email threads with information loss and no centralized view",
    "Desktop Tools (MS Project): Standalone tools without real-time collaboration or web access",
    "Physical Whiteboards: Sticky notes and printed Gantt charts that cannot be shared digitally",
    "No Authentication: File sharing with basic permissions — no role-based access control",
    "Manual Calculations: Progress percentages and resource utilization computed by hand",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(17), Inches(7.5),
                   existing_items, font_size=26, bold_prefix=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7: PROPOSED SYSTEM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Proposed System")

proposed_text = (
    "A full-stack MERN web application providing a centralized, real-time platform "
    "for managing projects, tasks, resources, and timelines — replacing manual workflows "
    "with automated, interactive features.\n\n"
)
add_body_text(slide, proposed_text, font_size=SMALL_SIZE, top=Inches(2.0), height=Inches(1.5))

proposed_items = [
    "Dashboard: Real-time KPI cards, Chart.js visualizations, activity feed and upcoming deadlines",
    "Project Management: Full CRUD operations, search/filter/sort/paginate, auto-delay detection",
    "Task Management: Dependencies and milestones, resource assignment, progress tracking",
    "Gantt Chart: Drag-and-drop scheduling with dhtmlx-gantt library, zoom levels and dependency arrows",
    "Resource Management: Workload utilization tracking, skills and availability, auto-computed metrics",
    "Authentication: JWT token-based auth, bcrypt password hashing, OTP-based password reset via email",
    "Email Service: Nodemailer with Ethereal test accounts and W3.css HTML email templates",
]
add_bullet_textbox(slide, Inches(0.8), Inches(3.6), Inches(17), Inches(6),
                   proposed_items, font_size=24, bold_prefix=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8: PROPOSED VS EXISTING
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Proposed System Vs Existing System")

comp_data = [
    ["Feature", "Existing System", "Proposed System"],
    ["Data Storage", "Scattered files & spreadsheets", "Centralized MongoDB database"],
    ["Real-Time Updates", "None", "Instant via REST API"],
    ["Gantt Charts", "Static (Excel)", "Interactive drag-and-drop"],
    ["Task Dependencies", "Not tracked", "Modeled & visualized"],
    ["Resource Tracking", "Manual calculation", "Auto-computed utilization"],
    ["Authentication", "File sharing permissions", "JWT + OTP-based reset"],
    ["Dashboard / KPIs", "Manual reports", "Auto-generated real-time charts"],
    ["Email Notifications", "None", "Nodemailer + W3.css templates"],
    ["Mobile Access", "Limited", "Fully responsive design"],
    ["Scalability", "Poor (file-based)", "High (MERN stack)"],
]
add_table(slide, Inches(1), Inches(2.0), Inches(16.5), Inches(7.5), comp_data,
          col_widths=[Inches(3.5), Inches(6.5), Inches(6.5)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 9: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "SYSTEM ARCHITECTURE")

# Three tiers
tiers = [
    ("PRESENTATION TIER (Frontend)", 
     "React.js 18 SPA  •  React Router v6  •  Chart.js  •  dhtmlx-gantt\n"
     "Context API (State)  •  Axios HTTP Client  •  Port 3000",
     BG_BLUE, Inches(2.0)),
    ("APPLICATION TIER (Backend)",
     "Node.js + Express.js (Port 5000)  •  JWT Auth Middleware  •  CORS\n"
     "Controllers: Auth | Project | Task | Resource | Dashboard  •  Nodemailer Email",
     BG_BLUE, Inches(4.5)),
    ("DATA TIER (Database)",
     "MongoDB Atlas (Cloud)  •  Mongoose ODM\n"
     "Collections: Users | Projects | Tasks | Resources",
     BG_BLUE, Inches(7.0)),
]

for title, desc, color, top in tiers:
    add_shape(slide, Inches(2.5), top, Inches(13.5), Inches(1.8), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(3), top + Inches(0.1), Inches(12.5), Inches(0.5),
                title, font_size=24, font_color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Times New Roman")
    add_textbox(slide, Inches(3), top + Inches(0.65), Inches(12.5), Inches(1.0),
                desc, font_size=18, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
                font_name="Times New Roman")

# Arrows between tiers
add_textbox(slide, Inches(8.5), Inches(3.85), Inches(2), Inches(0.5),
            "⬇ REST API / JSON + JWT ⬇", font_size=14, font_color=MEDIUM_BLUE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.5), Inches(6.35), Inches(2), Inches(0.5),
            "⬇ Mongoose ODM ⬇", font_size=14, font_color=MEDIUM_BLUE,
            alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10: USE CASE DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "USE CASE DIAGRAM")

# Actor column
add_textbox(slide, Inches(0.5), Inches(2.0), Inches(3.5), Inches(0.4),
            "Actors:", font_size=24, font_color=DARK_BLUE, bold=True)
actors = [
    "Project Manager — Full access",
    "Team Member — View + update tasks",
    "Admin — System administration",
]
add_bullet_textbox(slide, Inches(0.5), Inches(2.6), Inches(3.8), Inches(2.5),
                   actors, font_size=18, bullet="👤")

# Use case groups
uc_groups = [
    ("Authentication", ["Register", "Login / Logout", "Forgot Password (OTP)", "Reset Password"], Inches(4.5), Inches(2.0)),
    ("Project Management", ["Create / Edit / Delete Project", "View Projects & Details", "Search / Filter / Sort", "Track Progress"], Inches(9.5), Inches(2.0)),
    ("Task Management", ["Create / Update / Delete Task", "Assign Resource", "Set Dependency", "Mark Milestone"], Inches(4.5), Inches(5.5)),
    ("Resource & Dashboard", ["Add / Edit Resources", "View Utilization", "Dashboard KPIs & Charts", "Gantt Chart Visualization"], Inches(9.5), Inches(5.5)),
]

for title, cases, left, top in uc_groups:
    box_h = Inches(0.35 + len(cases) * 0.35)
    add_shape(slide, left, top, Inches(5.5), box_h, BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.05), Inches(5.1), Inches(0.35),
                title, font_size=18, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, uc in enumerate(cases):
        add_textbox(slide, left + Inches(0.3), top + Inches(0.4 + j * 0.32), Inches(5), Inches(0.3),
                    f"○ {uc}", font_size=15, font_color=DARK_TEXT)

add_textbox(slide, Inches(1), Inches(9.0), Inches(16), Inches(0.4),
            "28+ Use Cases  •  3 Actor Types  •  7 Functional Modules including Authentication & Email",
            font_size=16, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 11: CLASS DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "CLASS DIAGRAM")

classes_info = [
    ("User", Inches(0.3), [
        "_id: ObjectId", "name: String", "email: String (unique)",
        "password: String (hashed)", "role: user | admin",
        "resetPasswordOTP: String", "resetPasswordOTPExpire: Date",
    ], ["matchPassword()", "getSignedJwtToken()", "getResetPasswordOTP()"]),
    ("Project", Inches(4.8), [
        "_id: ObjectId", "name: String", "description: String",
        "startDate/endDate: Date", "status: Enum(4 states)",
        "completionPercentage: Number", "priority: Enum(4 levels)",
        "manager: String", "budget: Number",
    ], ["daysRemaining()", "totalDuration()"]),
    ("Task", Inches(9.5), [
        "_id: ObjectId", "projectId: → Project",
        "title: String", "startDate/endDate: Date",
        "status: Enum(4)", "assignedResource: → Resource",
        "dependencyTaskId: → Task", "milestone: Boolean",
        "progress: Number", "priority: Enum(4)",
    ], ["duration()", "isOverdue()"]),
    ("Resource", Inches(14.2), [
        "_id: ObjectId", "name: String", "email: String",
        "role: String", "department: String",
        "skills: String[]", "availability: Boolean",
        "maxHoursPerWeek: Number", "currentWorkload: Number",
    ], ["utilizationPct()", "availabilityStatus()"]),
]

for name, left, attrs, methods in classes_info:
    top = Inches(1.8)
    total_lines = len(attrs) + len(methods)
    box_h = Inches(0.5 + total_lines * 0.28)
    add_shape(slide, left, top, Inches(4.2), box_h, WHITE)

    # Header
    add_shape(slide, left, top, Inches(4.2), Inches(0.4), DARK_BLUE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.02), Inches(4), Inches(0.35),
                name, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    y = top + Inches(0.45)
    for attr in attrs:
        add_textbox(slide, left + Inches(0.15), y, Inches(3.9), Inches(0.26),
                    f"- {attr}", font_size=11, font_color=DARK_GRAY, font_name="Consolas")
        y += Inches(0.25)

    add_shape(slide, left + Inches(0.1), y, Inches(4), Inches(0.02), MEDIUM_BLUE)
    y += Inches(0.08)

    for method in methods:
        add_textbox(slide, left + Inches(0.15), y, Inches(3.9), Inches(0.26),
                    f"+ {method}", font_size=11, font_color=MEDIUM_BLUE, font_name="Consolas")
        y += Inches(0.25)

# Relationships
add_textbox(slide, Inches(0.3), Inches(8.5), Inches(18), Inches(0.5),
            "Relationships:   Project ──[1:*]──► Task      Task ──[*:0..1]──► Task (dependency)      Task ──[*:0..1]──► Resource",
            font_size=15, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12: SEQUENCE DIAGRAM — LOGIN
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "SEQUENCE DIAGRAM — User Login")

steps_login = [
    "1. User enters email and password on Login page",
    "2. React App calls AuthContext.login(credentials)",
    "3. Axios sends POST /api/auth/login with JSON body",
    "4. Express routes to authController.login()",
    "5. Controller queries MongoDB: User.findOne({email}).select('+password')",
    "6. Database returns user document",
    "7. Controller calls matchPassword() → getSignedJwtToken()",
    "8. Server responds: {success: true, token: JWT, user: {...}}",
    "9. Axios returns 200 OK with response data",
    "10. AuthContext stores token in localStorage, sets user state",
    "11. React navigates to Dashboard — user is authenticated",
]
add_bullet_textbox(slide, Inches(1), Inches(2.2), Inches(17), Inches(7),
                   steps_login, font_size=22, bullet="→")


# ════════════════════════════════════════════════════════════════════
# SLIDE 13: SEQUENCE DIAGRAM — PASSWORD RESET (OTP)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "SEQUENCE DIAGRAM — OTP Password Reset")

steps_otp = [
    "1. User clicks 'Forgot Password' and enters email",
    "2. POST /api/auth/forgotpassword — generates 6-digit OTP",
    "3. OTP is hashed (SHA-256) and stored in User document with 10-min expiry",
    "4. Plain-text OTP is sent via Nodemailer (Ethereal test or production SMTP)",
    "5. Server responds with success message (+ preview URL in dev mode)",
    "6. User enters 6-digit OTP on the verification step",
    "7. POST /api/auth/verifyotp — compares hashed OTP, checks expiry",
    "8. If valid: generates a short-lived reset token, clears OTP fields",
    "9. User enters new password with strength validation",
    "10. PUT /api/auth/resetpassword/:token — sets new password, clears token",
    "11. Server auto-logs in user and returns JWT — redirect to Dashboard",
]
add_bullet_textbox(slide, Inches(1), Inches(2.2), Inches(17), Inches(7),
                   steps_otp, font_size=22, bullet="→")


# ════════════════════════════════════════════════════════════════════
# SLIDE 14: ACTIVITY DIAGRAM — PROJECT LIFECYCLE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "ACTIVITY DIAGRAM — Project Lifecycle")

act_steps = [
    ("● Start", DARK_BLUE),
    ("Manager Creates Project", MEDIUM_BLUE),
    ("◇ Valid Data?", ORANGE),
    ("Save to DB (Status: Not Started)", ACCENT_GREEN),
    ("Add Tasks to Project", MEDIUM_BLUE),
    ("Assign Resources to Tasks", MEDIUM_BLUE),
    ("Project Status → In Progress", LIGHT_BLUE),
]
act_steps_right = [
    ("◇ All Tasks Completed?", ORANGE),
    ("◇ Past End Date?", ORANGE),
    ("Status → Completed (100%)", ACCENT_GREEN),
    ("Status → Delayed → Reschedule", ACCENT_RED),
    ("● End", DARK_BLUE),
]

for i, (text, color) in enumerate(act_steps):
    top = Inches(1.8) + i * Inches(0.9)
    add_shape(slide, Inches(1.5), top, Inches(6.5), Inches(0.6), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(1.7), top + Inches(0.08), Inches(6.1), Inches(0.45),
                text, font_size=17, font_color=WHITE, bold=True)

for i, (text, color) in enumerate(act_steps_right):
    top = Inches(1.8) + i * Inches(0.9)
    add_shape(slide, Inches(10), top, Inches(7), Inches(0.6), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(10.2), top + Inches(0.08), Inches(6.6), Inches(0.45),
                text, font_size=17, font_color=WHITE, bold=True)

add_textbox(slide, Inches(1), Inches(8.8), Inches(16), Inches(0.6),
            "Flow: Create → Add Tasks → Assign Resources → Track → Complete or Reschedule if Delayed",
            font_size=16, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15: ER DIAGRAM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "ENTITY-RELATIONSHIP (ER) DIAGRAM")

entities = [
    ("USER", Inches(0.3), [
        "PK _id: ObjectId", "name: String", "email: String (unique)",
        "password: String", "role: user | admin",
        "resetPasswordOTP: String", "resetPasswordToken: String",
    ]),
    ("PROJECT", Inches(4.8), [
        "PK _id: ObjectId", "name: String", "description: String",
        "startDate / endDate: Date", "status: Enum(4)",
        "completionPercentage: Number", "priority: Enum(4)",
        "manager: String", "budget: Number",
    ]),
    ("TASK", Inches(9.5), [
        "PK _id: ObjectId", "FK projectId → PROJECT",
        "title: String", "startDate / endDate: Date",
        "status: Enum(4)", "FK dependencyTaskId → TASK",
        "FK assignedResource → RESOURCE",
        "milestone: Boolean", "progress: Number",
    ]),
    ("RESOURCE", Inches(14.2), [
        "PK _id: ObjectId", "name: String", "email: String",
        "role: String", "department: String",
        "availability: Boolean", "skills: String[]",
        "maxHoursPerWeek: Number",
    ]),
]

for name, left, fields in entities:
    top = Inches(1.8)
    h = Inches(0.3 * len(fields) + 0.6)
    add_shape(slide, left, top, Inches(4.2), h, WHITE)
    add_shape(slide, left, top, Inches(4.2), Inches(0.4), DARK_BLUE)
    add_textbox(slide, left, top + Inches(0.02), Inches(4.2), Inches(0.35),
                name, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, field in enumerate(fields):
        fc = MEDIUM_BLUE if field.startswith("PK") or field.startswith("FK") else DARK_GRAY
        add_textbox(slide, left + Inches(0.1), top + Inches(0.48 + j * 0.28),
                    Inches(4), Inches(0.26),
                    field, font_size=12, font_color=fc, font_name="Consolas")

add_textbox(slide, Inches(0.5), Inches(7.0), Inches(17), Inches(1.0),
            "Relationships:\n"
            "PROJECT ─── [1 : Many] ───► TASK    |    "
            "TASK ─── [Many : 0..1] ───► TASK (dependency)    |    "
            "TASK ─── [Many : 0..1] ───► RESOURCE (assignment)",
            font_size=15, font_color=DARK_BLUE, line_spacing=1.5)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16: MODULES & KEY FEATURES
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "MODULES & KEY FEATURES")

modules = [
    ("📊 Dashboard", ["Real-time KPI cards (6 metrics)", "Task status doughnut chart",
                        "Project progress bar chart", "Recent activities & deadlines"]),
    ("📁 Projects", ["Create/Edit/Delete projects", "Search, filter, sort, paginate",
                      "Priority & status badges", "Auto-delay detection"]),
    ("✅ Tasks", ["CRUD per project + global view", "Dependency & milestone flags",
                   "Resource assignment", "Progress % tracking"]),
    ("📅 Gantt Chart", ["dhtmlx-gantt library", "Drag-and-drop scheduling",
                         "Dependency arrows", "Zoom: Hours/Days/Weeks/Months"]),
    ("👥 Resources", ["Team member profiles", "Utilization % auto-computed",
                       "Availability status", "Skills & department info"]),
    ("🔐 Auth & Email", ["JWT + OTP password reset", "Nodemailer email service",
                          "bcrypt password hashing", "Register/Login/ForgotPwd"]),
]

for i, (title, features) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(6.2)
    top = Inches(2.0) + row * Inches(3.8)

    add_shape(slide, left, top, Inches(5.8), Inches(3.2), BG_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(0.4),
                title, font_size=20, font_color=DARK_BLUE, bold=True)
    add_bullet_textbox(slide, left + Inches(0.3), top + Inches(0.7), Inches(5.2), Inches(2.3),
                       features, font_size=16, font_color=DARK_GRAY, bullet="▸")


# ════════════════════════════════════════════════════════════════════
# SLIDE 17: API ENDPOINTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "API ENDPOINTS")

api_data = [
    ["Method", "Endpoint", "Description"],
    ["POST", "/api/auth/register", "Register new user"],
    ["POST", "/api/auth/login", "Login & get JWT token"],
    ["POST", "/api/auth/forgotpassword", "Send OTP to email"],
    ["POST", "/api/auth/verifyotp", "Verify 6-digit OTP"],
    ["PUT", "/api/auth/resetpassword/:token", "Reset password with token"],
    ["GET", "/api/auth/me", "Get current user (protected)"],
    ["GET", "/api/projects", "List all projects"],
    ["POST", "/api/projects", "Create project"],
    ["PUT", "/api/projects/:id", "Update project"],
    ["DELETE", "/api/projects/:id", "Delete project"],
    ["GET", "/api/tasks/all", "Get all tasks"],
    ["POST", "/api/tasks", "Create task"],
    ["PUT", "/api/tasks/:id", "Update task"],
    ["GET", "/api/tasks/gantt/:id", "Get Gantt chart data"],
    ["GET", "/api/resources", "List all resources"],
    ["GET", "/api/dashboard/overview", "Dashboard statistics"],
]
add_table(slide, Inches(1.5), Inches(2.0), Inches(15.5), Inches(7.0), api_data,
          col_widths=[Inches(2), Inches(6), Inches(7.5)])


# ════════════════════════════════════════════════════════════════════
# SLIDE 18: TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "TECHNOLOGY STACK")

# Frontend table
add_textbox(slide, Inches(0.5), Inches(1.8), Inches(8), Inches(0.4),
            "Frontend", font_size=24, font_color=DARK_BLUE, bold=True)
fe_stack = [
    ["Technology", "Purpose"],
    ["React.js 18", "UI framework with hooks & context"],
    ["React Router v6", "Client-side routing"],
    ["Chart.js", "Doughnut, bar & line charts"],
    ["dhtmlx-gantt", "Interactive Gantt chart"],
    ["Axios", "HTTP client with interceptors"],
    ["react-hot-toast", "Toast notifications"],
    ["React Icons", "Icon library"],
]
add_table(slide, Inches(0.5), Inches(2.3), Inches(8.5), Inches(3.0), fe_stack,
          col_widths=[Inches(3), Inches(5.5)])

# Backend table
add_textbox(slide, Inches(10), Inches(1.8), Inches(8), Inches(0.4),
            "Backend", font_size=24, font_color=DARK_BLUE, bold=True)
be_stack = [
    ["Technology", "Purpose"],
    ["Node.js", "JavaScript runtime"],
    ["Express.js 4", "Web framework"],
    ["MongoDB + Mongoose", "NoSQL database + ODM"],
    ["JWT (jsonwebtoken)", "Authentication tokens"],
    ["bcryptjs", "Password hashing"],
    ["Nodemailer", "Email sending (OTP, W3.css templates)"],
    ["Ethereal Email", "Auto test accounts (development)"],
]
add_table(slide, Inches(10), Inches(2.3), Inches(8), Inches(3.0), be_stack,
          col_widths=[Inches(3.5), Inches(4.5)])

# Dev tools
add_textbox(slide, Inches(0.5), Inches(6.0), Inches(8), Inches(0.4),
            "Development Tools", font_size=24, font_color=DARK_BLUE, bold=True)
dev_items = [
    "concurrently — Run frontend + backend simultaneously",
    "nodemon — Auto-restart server on changes",
    "VS Code — IDE with extensions",
    "Postman — API testing",
    "Git — Version control",
]
add_bullet_textbox(slide, Inches(0.5), Inches(6.5), Inches(17), Inches(2.5),
                   dev_items, font_size=18, bullet="▸")


# ════════════════════════════════════════════════════════════════════
# SLIDE 19: AUTHENTICATION SYSTEM
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "AUTHENTICATION SYSTEM")

auth_text = (
    "The application implements a complete authentication system with the following features:\n\n"
)
add_body_text(slide, auth_text, font_size=SMALL_SIZE, top=Inches(2.0), height=Inches(1))

auth_features = [
    "Register: User registration with name, email, password validation and bcrypt hashing (salt rounds: 10)",
    "Login: Email/password authentication returning JWT token stored in localStorage",
    "Protected Routes: PrivateRoute component checks auth state — redirects to login if unauthenticated",
    "Forgot Password: 3-step OTP flow — Enter email → Receive 6-digit OTP → Set new password",
    "OTP Email: Nodemailer sends W3.css-styled HTML emails via Ethereal (dev) or production SMTP",
    "OTP Verification: SHA-256 hashed OTP stored in MongoDB, 10-minute expiry, single-use",
    "Password Strength: Real-time client-side validation (length, uppercase, lowercase, digit, special char)",
    "Auto Sign-In: After successful password reset, user receives JWT and is redirected to Dashboard",
]
add_bullet_textbox(slide, Inches(0.8), Inches(3.2), Inches(17), Inches(6),
                   auth_features, font_size=20, bold_prefix=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 20: EMAIL SERVICE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "EMAIL SERVICE — Nodemailer + W3.css")

email_features = [
    "Auto-Initialization: On server start, tries production SMTP credentials first — automatically falls back to Ethereal test account if credentials are missing or invalid",
    "Ethereal Email: Nodemailer's free test service — auto-creates disposable SMTP accounts, captures all emails, provides preview URLs to view emails in browser",
    "W3.css HTML Templates: Beautiful responsive email templates using W3.css framework with gradient headers, OTP code display, action buttons, and security notices",
    "OTP Email: Sends 6-digit verification code with styled template — large code display, 10-minute expiry notice, security warning",
    "Welcome Email: Sent on registration with dashboard link and feature highlights",
    "Password Change Confirmation: Sent after successful password reset with timestamp and sign-in link",
    "Three Modes: Production (real SMTP) → Ethereal (test capture) → Console (fallback logging)",
    "Preview URLs: In development mode, Ethereal preview links are returned to the frontend so developers can view sent emails",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(17), Inches(7),
                   email_features, font_size=19, bold_prefix=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 21: RESULTS — DASHBOARD
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Results — Dashboard View")

add_shape(slide, Inches(2), Inches(2.0), Inches(14.5), Inches(7), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(4), Inches(4.5), Inches(10.5), Inches(1.5),
            "[ Dashboard Screenshot ]\n\nReal-time KPI cards, task distribution chart,\nproject progress bars, recent activities",
            font_size=20, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 22: RESULTS — PROJECTS & GANTT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Results — Projects & Gantt Chart")

# Left placeholder
add_shape(slide, Inches(0.5), Inches(2.0), Inches(8.5), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(1.5), Inches(4.0), Inches(6.5), Inches(1.5),
            "[ Projects Page Screenshot ]\n\nProject cards with status badges,\nsearch & filter controls",
            font_size=18, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Right placeholder
add_shape(slide, Inches(9.5), Inches(2.0), Inches(9), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(10.5), Inches(4.0), Inches(7), Inches(1.5),
            "[ Gantt Chart Screenshot ]\n\nInteractive timeline with drag-and-drop,\ntask dependencies and milestones",
            font_size=18, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Labels
add_textbox(slide, Inches(0.5), Inches(7.7), Inches(8.5), Inches(0.4),
            "Projects Page", font_size=20, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.5), Inches(7.7), Inches(9), Inches(0.4),
            "Gantt Chart Page", font_size=20, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 23: RESULTS — AUTH PAGES
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Results — Authentication Pages")

# Three placeholders
labels = ["Login Page", "Forgot Password (OTP Flow)", "Password Reset"]
for i, label in enumerate(labels):
    left = Inches(0.5) + i * Inches(6.2)
    add_shape(slide, left, Inches(2.0), Inches(5.8), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left + Inches(0.5), Inches(4.2), Inches(4.8), Inches(1),
                f"[ {label} Screenshot ]", font_size=18, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(7.7), Inches(5.8), Inches(0.4),
                label, font_size=20, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 24: FUTURE ENHANCEMENTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "Future Enhancements")

enhancements = [
    "Email Notifications: Deadline reminders, task assignment alerts, weekly digest reports using production SMTP",
    "Advanced Reporting: PDF/Excel export, custom report builder, budget vs actual analysis",
    "Real-Time Collaboration: Comments on tasks, file attachments, real-time team chat via WebSocket",
    "Mobile Application: React Native version with offline support and push notifications",
    "AI-Powered Features: Automated task estimation, risk prediction, resource optimization using ML models",
    "Advanced Analytics: Burndown charts, velocity tracking, predictive timelines using historical data",
    "Role-Based Access Control: Granular permissions — Admin, Manager, Member roles with different access levels",
    "Multi-Language Support: Internationalization (i18n) for supporting multiple languages",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(17), Inches(7),
                   enhancements, font_size=24, bold_prefix=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 25: CONCLUSION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, "CONCLUSION")

conclusion = (
    "We have successfully developed a comprehensive, full-stack Project Management Dashboard "
    "using the MERN stack (MongoDB, Express.js, React.js, Node.js) that addresses the key "
    "challenges of modern project management.\n\n"
    "The system provides interactive Gantt chart visualization with drag-and-drop scheduling, "
    "real-time dashboard analytics with KPI cards and Chart.js visualizations, complete CRUD "
    "operations for projects, tasks, and resources, and a robust authentication system with "
    "JWT tokens, OTP-based password reset, and Nodemailer email integration.\n\n"
    "The application demonstrates practical implementation of full-stack web development, "
    "RESTful API design, database modeling with Mongoose, state management with React Context API, "
    "and responsive UI design. The modular architecture ensures easy maintenance and scalability "
    "for future enhancements including real-time collaboration, mobile apps, and AI-driven features."
)
add_body_text(slide, conclusion, font_size=SMALL_SIZE, line_spacing=1.45)


# ════════════════════════════════════════════════════════════════════
# SLIDE 26: THANK YOU
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Decorative circles
add_decorative_circle(slide, Inches(17), Inches(-0.5), Inches(2.5), BG_BLUE)
add_decorative_circle(slide, Inches(0), Inches(8.5), Inches(2), LIGHT_BLUE)

add_textbox(slide, Inches(2), Inches(2.5), Inches(16), Inches(1.5),
            "THANK YOU", font_size=60, font_color=DARK_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Times New Roman")

add_textbox(slide, Inches(2), Inches(4.5), Inches(16), Inches(1.0),
            "Project Management Dashboard\nwith Interactive Gantt Chart Visualization System",
            font_size=28, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
            font_name="Times New Roman")

add_textbox(slide, Inches(3), Inches(6.0), Inches(14), Inches(0.5),
            "Tech Stack: React.js | Node.js | Express.js | MongoDB | Chart.js | DHTMLX Gantt",
            font_size=20, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(5), Inches(8.0), Inches(10), Inches(0.8),
            "Q & A", font_size=44, font_color=DARK_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Times New Roman")


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"✅ Presentation saved to: {OUTPUT_PATH}")
print(f"📊 Total slides: {len(prs.slides)}")
